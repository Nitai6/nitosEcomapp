#!/usr/bin/env python3
"""
Generate a styled PowerPoint presentation from a JSON outline.

Usage:
    python3 generate_presentation.py \
        --outline presentations/my-talk-assets/outline.json \
        --output presentations/my-talk.pptx \
        --mode normal

Modes: normal, camera-left, camera-right

Environment:
    Requires python-pptx: pip install python-pptx
"""

import argparse
import copy
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


# ── Design System ──────────────────────────────────────────────

COLORS = {
    "background": RGBColor(0x0A, 0x0B, 0x12),
    "accent": RGBColor(0x00, 0xCB, 0x5A),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "light_gray": RGBColor(0xA0, 0xA0, 0xA0),
    "dark_gray": RGBColor(0x1A, 0x1B, 0x26),
    "card_border": RGBColor(0x2A, 0x2B, 0x36),
}

FONTS = {
    "heading": "Heading Now Trial 04",       # Set bold=True for Bold weight
    "body": "Heading Now Trial 04",          # Regular weight (bold=False)
    "serif": "IBM Plex Serif",               # Medium Italic for accents
    "comic": "Steel City Comic",             # Playful accent font
}

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Margins
MARGIN = Inches(0.8)
MARGIN_TOP = Inches(0.6)


# ── Layout Bounds ──────────────────────────────────────────────

def get_content_bounds(mode):
    """Return (left, top, width, height) for content area based on mode."""
    if mode == "camera-left":
        content_left = SLIDE_WIDTH * 0.4
        content_width = SLIDE_WIDTH * 0.6 - MARGIN
        return (content_left + Inches(0.3), MARGIN_TOP, content_width, SLIDE_HEIGHT - MARGIN_TOP - MARGIN)
    elif mode == "camera-right":
        content_width = SLIDE_WIDTH * 0.6 - MARGIN
        return (MARGIN, MARGIN_TOP, content_width, SLIDE_HEIGHT - MARGIN_TOP - MARGIN)
    else:
        return (MARGIN, MARGIN_TOP, SLIDE_WIDTH - MARGIN * 2, SLIDE_HEIGHT - MARGIN_TOP - MARGIN)


def get_image_bounds(mode):
    """Return (left, top, width, height) for full-bleed images based on mode."""
    if mode == "camera-left":
        left = SLIDE_WIDTH * 0.4
        return (left, 0, SLIDE_WIDTH - left, SLIDE_HEIGHT)
    elif mode == "camera-right":
        return (0, 0, SLIDE_WIDTH * 0.6, SLIDE_HEIGHT)
    else:
        return (0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)


# ── Helpers ────────────────────────────────────────────────────

def apply_background(slide):
    """Set dark background on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["background"]


def add_text_box(slide, left, top, width, height, text, font_name, font_size,
                 color=None, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP):
    """Add a styled single-paragraph text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None

    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.alignment = alignment

    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color or COLORS["white"]
    run.font.bold = bold
    run.font.italic = italic

    return txBox


def parse_body_text(body_text):
    """Parse a body string into paragraph dicts. Lines starting with '- ' or '* ' become bullets."""
    paragraphs = []
    for line in body_text.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("- ") or line.startswith("* "):
            paragraphs.append({"text": line[2:], "bullet": True})
        else:
            paragraphs.append({"text": line, "bullet": False})
    return paragraphs


def add_rich_text_box(slide, left, top, width, height, paragraphs_data,
                      default_font, default_size, default_color=None,
                      alignment=PP_ALIGN.LEFT):
    """Add a text box with multiple paragraphs, supporting bullets.

    paragraphs_data: list of dicts with keys:
        - text (str): paragraph text
        - bullet (bool): render as bullet item (default False)
        - bold (bool): override bold (default False)
        - italic (bool): override italic (default False)
        - font (str): override font name
        - size (int): override font size in pt
        - color (RGBColor): override color
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    color = default_color or COLORS["white"]

    for i, pdata in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = alignment

        # Set paragraph-level properties via XML
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()

        if pdata.get("bullet"):
            # Bullet character
            buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
            pPr.append(buChar)

            # Bullet color (green accent)
            buClr = pPr.makeelement(qn("a:buClr"), {})
            srgbClr = buClr.makeelement(qn("a:srgbClr"), {"val": "00CB5A"})
            buClr.append(srgbClr)
            pPr.append(buClr)

            # Bullet size (100% of text)
            buSzPct = pPr.makeelement(qn("a:buSzPct"), {"val": "100000"})
            pPr.append(buSzPct)

            # Indentation
            pPr.set("marL", str(int(Inches(0.4))))
            pPr.set("indent", str(int(Inches(-0.25))))

            # Space after bullet
            spcAft = pPr.makeelement(qn("a:spcAft"), {})
            spcPts = spcAft.makeelement(qn("a:spcPts"), {"val": "600"})
            spcAft.append(spcPts)
            pPr.append(spcAft)
        else:
            # Suppress inherited bullets
            buNone = pPr.makeelement(qn("a:buNone"), {})
            pPr.append(buNone)

            # Space after regular paragraph
            spcAft = pPr.makeelement(qn("a:spcAft"), {})
            spcPts = spcAft.makeelement(qn("a:spcPts"), {"val": "400"})
            spcAft.append(spcPts)
            pPr.append(spcAft)

        # Add run with text
        run = p.add_run()
        run.text = pdata.get("text", "")

        # Font properties
        font = run.font
        font.name = pdata.get("font", default_font)
        font.size = Pt(pdata.get("size", default_size))
        font.color.rgb = pdata.get("color", color)
        font.bold = pdata.get("bold", False)
        font.italic = pdata.get("italic", False)

    return txBox


def add_accent_line(slide, left, top, width, height=None):
    """Add a decorative accent line (thin filled rectangle)."""
    if height is None:
        height = Pt(3)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["accent"]
    shape.line.fill.background()
    return shape


def add_number_circle(slide, left, top, number, diameter=None):
    """Add a filled accent circle with a number inside."""
    if diameter is None:
        diameter = Inches(0.55)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["accent"]
    shape.line.fill.background()

    # Add number text inside
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{number:02d}" if isinstance(number, int) else str(number)
    run.font.name = FONTS["heading"]
    run.font.size = Pt(18)
    run.font.color.rgb = COLORS["background"]
    run.font.bold = True

    # Vertical center
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    return shape


def add_card_background(slide, left, top, width, height):
    """Add a subtle dark card background (rounded rectangle)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["dark_gray"]
    shape.line.color.rgb = COLORS["card_border"]
    shape.line.width = Pt(1)
    # Reduce corner rounding
    shape.adjustments[0] = 0.05
    return shape


def add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide.

    Note: python-pptx 1.0.2 notes_slide creation can corrupt PPTX files
    in some environments (Keynote, older PowerPoint). Speaker notes are
    stored in the outline JSON instead and skipped during generation.
    """
    pass


def fit_image(slide, image_path, bounds):
    """Add an image that fits within bounds, maintaining aspect ratio."""
    left, top, max_w, max_h = bounds
    try:
        from PIL import Image as PILImage
        with PILImage.open(image_path) as img:
            img_w, img_h = img.size
    except Exception:
        slide.shapes.add_picture(str(image_path), left, top, max_w, max_h)
        return

    scale_w = max_w / Emu(int(img_w * 914400 / 96))
    scale_h = max_h / Emu(int(img_h * 914400 / 96))
    scale = min(scale_w, scale_h, 1.5)

    pic_w = int(img_w * 914400 / 96 * scale)
    pic_h = int(img_h * 914400 / 96 * scale)

    pic_left = left + (max_w - pic_w) // 2
    pic_top = top + (max_h - pic_h) // 2

    slide.shapes.add_picture(str(image_path), pic_left, pic_top, pic_w, pic_h)


# ── Slide Renderers ────────────────────────────────────────────

def render_title(slide, data, bounds):
    """Big centered title + optional subtitle."""
    left, top, width, height = bounds
    heading = data.get("heading", "")
    subtitle = data.get("subtitle", "")

    title_top = top + height * 0.25 if subtitle else top + height * 0.3

    # Title
    add_text_box(
        slide, left, int(title_top), width, int(height * 0.35),
        heading, FONTS["heading"], 54,
        color=COLORS["white"], bold=True, alignment=PP_ALIGN.CENTER,
    )

    # Subtitle
    if subtitle:
        add_text_box(
            slide, left, int(title_top + Inches(1.8)), width, int(height * 0.2),
            subtitle, FONTS["serif"], 22,
            color=COLORS["light_gray"], italic=True, alignment=PP_ALIGN.CENTER,
        )


def render_content(slide, data, bounds):
    """Heading + rich body text with optional bullet support."""
    left, top, width, height = bounds
    heading = data.get("heading", "")
    body = data.get("body", "")
    bullets = data.get("bullets", [])

    if heading:
        add_text_box(
            slide, left, int(top + Inches(0.5)), width, int(Inches(1.2)),
            heading, FONTS["heading"], 36,
            color=COLORS["accent"], bold=True, alignment=PP_ALIGN.LEFT,
        )

    body_top = int(top + Inches(2.0)) if heading else int(top + Inches(0.5))
    body_height = int(height - Inches(2.5)) if heading else int(height - Inches(1.0))

    # Build paragraphs
    paragraphs = []
    if body:
        paragraphs.extend(parse_body_text(body))
    if bullets:
        for b in bullets:
            paragraphs.append({"text": b, "bullet": True})

    if paragraphs:
        add_rich_text_box(
            slide, left, body_top, width, body_height,
            paragraphs, FONTS["body"], 26,
            default_color=COLORS["white"], alignment=PP_ALIGN.LEFT,
        )


def render_content_media(slide, data, bounds, image_bounds):
    """Text + GIF/image side by side."""
    left, top, width, height = bounds
    heading = data.get("heading", "")
    body = data.get("body", "")
    bullets = data.get("bullets", [])
    image_path = data.get("image_path", "")
    media_side = data.get("media_side", "right")
    caption = data.get("caption", "")

    gap = Inches(0.3)
    text_pct = 0.55
    media_pct = 0.45

    if media_side == "right":
        text_left = left
        text_width = int(width * text_pct - gap)
        media_left = int(left + width * text_pct + gap / 2)
        media_width = int(width * media_pct - gap / 2)
    else:
        media_left = left
        media_width = int(width * media_pct - gap / 2)
        text_left = int(left + width * media_pct + gap / 2)
        text_width = int(width * text_pct - gap)

    # --- Text side ---
    if heading:
        add_text_box(
            slide, text_left, int(top + Inches(0.5)), text_width, int(Inches(1.2)),
            heading, FONTS["heading"], 32,
            color=COLORS["accent"], bold=True, alignment=PP_ALIGN.LEFT,
        )

    txt_body_top = int(top + Inches(1.8)) if heading else int(top + Inches(0.5))
    paragraphs = []
    if body:
        paragraphs.extend(parse_body_text(body))
    if bullets:
        for b in bullets:
            paragraphs.append({"text": b, "bullet": True})

    if paragraphs:
        add_rich_text_box(
            slide, text_left, txt_body_top, text_width, int(height - Inches(2.3)),
            paragraphs, FONTS["body"], 22,
            default_color=COLORS["white"], alignment=PP_ALIGN.LEFT,
        )

    # --- Media side ---
    media_top = int(top + Inches(0.3))
    media_height = int(height - Inches(0.6))

    # Card background
    add_card_background(slide, media_left, media_top, media_width, media_height)

    # Image/GIF
    padding = Inches(0.15)
    caption_reserve = Inches(0.4) if caption else 0
    if image_path and Path(image_path).exists():
        fit_image(slide, image_path, (
            int(media_left + padding),
            int(media_top + padding),
            int(media_width - padding * 2),
            int(media_height - padding * 2 - caption_reserve),
        ))
    else:
        add_text_box(
            slide, media_left, int(media_top + media_height * 0.4), media_width, int(Inches(1)),
            f"[Media: {image_path}]", FONTS["body"], 16,
            color=COLORS["light_gray"], alignment=PP_ALIGN.CENTER,
        )

    # Caption
    if caption:
        add_text_box(
            slide, media_left, int(media_top + media_height - Inches(0.35)),
            media_width, int(Inches(0.3)),
            caption, FONTS["body"], 14,
            color=COLORS["light_gray"], alignment=PP_ALIGN.CENTER,
        )


def render_image(slide, data, bounds, image_bounds):
    """Full-bleed image with optional caption."""
    image_path = data.get("image_path", "")
    caption = data.get("caption", "")

    if image_path and Path(image_path).exists():
        fit_image(slide, image_path, image_bounds)
    else:
        left, top, width, height = bounds
        add_text_box(
            slide, left, int(top + height * 0.4), width, int(Inches(1)),
            f"[Image: {image_path}]", FONTS["body"], 20,
            color=COLORS["light_gray"], alignment=PP_ALIGN.CENTER,
        )

    if caption:
        left, top, width, height = bounds
        add_text_box(
            slide, left, int(top + height - Inches(0.8)), width, int(Inches(0.6)),
            caption, FONTS["body"], 16,
            color=COLORS["light_gray"], alignment=PP_ALIGN.CENTER,
        )


def render_gif(slide, data, bounds, image_bounds):
    """Embedded GIF with optional caption."""
    render_image(slide, data, bounds, image_bounds)


def render_section_divider(slide, data, bounds):
    """Section transition with optional number, accent line, heading, subtitle."""
    left, top, width, height = bounds
    heading = data.get("heading", "")
    subtitle = data.get("subtitle", "")
    section_number = data.get("section_number", None)

    center_x = int(left + width / 2)

    if section_number is not None:
        # Number circle
        circle_size = Inches(0.55)
        add_number_circle(
            slide,
            int(center_x - circle_size / 2),
            int(top + Inches(1.0)),
            section_number,
        )
        line_top = int(top + Inches(1.8))
    else:
        line_top = int(top + Inches(1.5))

    # Accent line
    line_width = Inches(2.5)
    add_accent_line(
        slide,
        int(center_x - line_width / 2),
        line_top,
        int(line_width),
    )

    # Heading
    heading_top = int(line_top + Inches(0.4))
    add_text_box(
        slide, left, heading_top, width, int(Inches(1.5)),
        heading, FONTS["heading"], 54,
        color=COLORS["accent"], bold=True, alignment=PP_ALIGN.CENTER,
    )

    # Subtitle
    if subtitle:
        add_text_box(
            slide, left, int(heading_top + Inches(1.2)), width, int(Inches(0.8)),
            subtitle, FONTS["serif"], 20,
            color=COLORS["light_gray"], italic=True, alignment=PP_ALIGN.CENTER,
        )


def render_quote(slide, data, bounds):
    """Styled quote with vertical accent bar."""
    left, top, width, height = bounds
    text = data.get("text", "")
    attribution = data.get("attribution", "")

    # Vertical accent bar
    bar_top = int(top + height * 0.25)
    bar_height = Inches(2.2)
    add_accent_line(
        slide,
        int(left + Inches(0.3)),
        bar_top,
        Pt(4),
        height=int(bar_height),
    )

    # Quote text
    add_text_box(
        slide, int(left + Inches(0.9)), bar_top, int(width - Inches(1.2)), int(Inches(1.8)),
        text, FONTS["serif"], 30,
        color=COLORS["white"], italic=True, alignment=PP_ALIGN.LEFT,
    )

    # Attribution
    if attribution:
        add_text_box(
            slide, int(left + Inches(0.9)), int(bar_top + Inches(2.0)),
            int(width - Inches(1.2)), int(Inches(0.5)),
            f"\u2014 {attribution}", FONTS["body"], 18,
            color=COLORS["light_gray"], alignment=PP_ALIGN.LEFT,
        )


def render_stats(slide, data, bounds):
    """Big impactful statistic with label and optional context."""
    left, top, width, height = bounds
    value = data.get("value", "")
    label = data.get("label", "")
    context = data.get("context", "")

    center_y = int(top + height * 0.25)

    # Big number
    add_text_box(
        slide, left, center_y, width, int(Inches(1.8)),
        value, FONTS["heading"], 96,
        color=COLORS["accent"], bold=True, alignment=PP_ALIGN.CENTER,
    )

    # Accent line below number
    line_width = Inches(2.0)
    add_accent_line(
        slide,
        int(left + width / 2 - line_width / 2),
        int(center_y + Inches(1.6)),
        int(line_width),
        height=Pt(2),
    )

    # Label
    add_text_box(
        slide, left, int(center_y + Inches(2.0)), width, int(Inches(1.0)),
        label, FONTS["body"], 28,
        color=COLORS["white"], alignment=PP_ALIGN.CENTER,
    )

    # Context
    if context:
        add_text_box(
            slide, left, int(center_y + Inches(2.8)), width, int(Inches(0.6)),
            context, FONTS["serif"], 18,
            color=COLORS["light_gray"], italic=True, alignment=PP_ALIGN.CENTER,
        )


def render_two_column(slide, data, bounds):
    """Side-by-side columns for comparisons."""
    left, top, width, height = bounds
    heading = data.get("heading", "")

    col_top = int(top + Inches(0.3))
    if heading:
        add_text_box(
            slide, left, col_top, width, int(Inches(1.0)),
            heading, FONTS["heading"], 32,
            color=COLORS["accent"], bold=True, alignment=PP_ALIGN.LEFT,
        )
        col_top = int(top + Inches(1.6))

    col_height = int(top + height - col_top - Inches(0.3))
    col_width = int(width * 0.47)
    right_col_left = int(left + width * 0.53)

    # Vertical divider
    divider_x = int(left + width * 0.5 - Pt(1))
    add_accent_line(
        slide, divider_x, int(col_top + Inches(0.2)),
        Pt(2), height=int(col_height - Inches(0.4)),
    )

    for side in ("left", "right"):
        col_left = left if side == "left" else right_col_left
        col_heading = data.get(f"{side}_heading", "")
        col_body = data.get(f"{side}_body", "")
        col_bullets = data.get(f"{side}_bullets", [])

        y_offset = col_top

        if col_heading:
            add_text_box(
                slide, col_left, y_offset, col_width, int(Inches(0.8)),
                col_heading, FONTS["heading"], 28,
                color=COLORS["white"], bold=True, alignment=PP_ALIGN.LEFT,
            )
            y_offset = int(y_offset + Inches(1.0))

        paragraphs = []
        if col_body:
            paragraphs.extend(parse_body_text(col_body))
        if col_bullets:
            for b in col_bullets:
                paragraphs.append({"text": b, "bullet": True})

        if paragraphs:
            add_rich_text_box(
                slide, col_left, y_offset, col_width, int(col_height - Inches(1.2)),
                paragraphs, FONTS["body"], 22,
                default_color=COLORS["white"], alignment=PP_ALIGN.LEFT,
            )


def render_accent(slide, data, bounds):
    """Playful text using Steel City Comic font."""
    left, top, width, height = bounds
    text = data.get("text", "")

    add_text_box(
        slide, left, int(top + height * 0.3), width, int(height * 0.4),
        text, FONTS["comic"], 48,
        color=COLORS["accent"], bold=True, alignment=PP_ALIGN.CENTER,
    )


def render_blank(slide, data, bounds):
    """Just the dark background. Breathing room."""
    pass


# ── Main ───────────────────────────────────────────────────────

RENDERERS = {
    "title": render_title,
    "content": render_content,
    "content_media": render_content_media,
    "image": render_image,
    "gif": render_gif,
    "section_divider": render_section_divider,
    "quote": render_quote,
    "accent": render_accent,
    "blank": render_blank,
    "stats": render_stats,
    "two_column": render_two_column,
}

# These renderers need image_bounds in addition to content_bounds
MEDIA_TYPES = {"image", "gif", "content_media"}


def parse_args():
    parser = argparse.ArgumentParser(description="Generate styled PPTX from JSON outline")
    parser.add_argument("--outline", required=True, help="Path to JSON outline file")
    parser.add_argument("--output", required=True, help="Output .pptx file path")
    parser.add_argument(
        "--mode", default="normal",
        choices=["normal", "camera-left", "camera-right"],
        help="Layout mode (default: normal)"
    )
    return parser.parse_args()


def main():
    args = parse_args()

    outline_path = Path(args.outline)
    if not outline_path.exists():
        print(f"Error: Outline file not found: {outline_path}", file=sys.stderr)
        sys.exit(1)

    with open(outline_path) as f:
        outline = json.load(f)

    mode = args.mode if args.mode != "normal" else outline.get("mode", "normal")

    content_bounds = get_content_bounds(mode)
    image_bounds = get_image_bounds(mode)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slides = outline.get("slides", [])
    if not slides:
        print("Error: No slides in outline", file=sys.stderr)
        sys.exit(1)

    blank_layout = prs.slide_layouts[6]

    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        apply_background(slide)

        slide_type = slide_data.get("type", "content")
        renderer = RENDERERS.get(slide_type)

        if not renderer:
            print(f"Warning: Unknown slide type '{slide_type}' at index {i}, skipping", file=sys.stderr)
            continue

        if slide_type in MEDIA_TYPES:
            renderer(slide, slide_data, content_bounds, image_bounds)
        else:
            renderer(slide, slide_data, content_bounds)

        add_speaker_notes(slide, slide_data.get("notes", ""))

    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(slides)}")
    print(f"Mode: {mode}")


if __name__ == "__main__":
    main()
